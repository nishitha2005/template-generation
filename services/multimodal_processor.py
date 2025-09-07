import os
import PyPDF2
import openpyxl
from pptx import Presentation
from docx import Document
from PIL import Image
import speech_recognition as sr
from pydub import AudioSegment
import cv2
import json
import re
from typing import Dict, List, Any

class MultimodalProcessor:
    """Process various file types and extract content with metadata"""
    
    def __init__(self):
        self.supported_types = {
            '.pdf': self._process_pdf,
            '.pptx': self._process_pptx,
            '.docx': self._process_docx,
            '.xlsx': self._process_xlsx,
            '.xls': self._process_xlsx,
            '.png': self._process_image,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.mp3': self._process_audio,
            '.wav': self._process_audio,
            '.mp4': self._process_video,
            '.avi': self._process_video
        }
    
    def get_file_type(self, filename: str) -> str:
        """Get file type based on extension"""
        ext = os.path.splitext(filename)[1].lower()
        return ext[1:] if ext in self.supported_types else 'unknown'
    
    def process_file(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Process file and extract content with metadata"""
        file_type = self.get_file_type(filename)
        
        if file_type not in [ext[1:] for ext in self.supported_types.keys()]:
            return {
                'type': 'unsupported',
                'content': '',
                'metadata': {'error': 'Unsupported file type'}
            }
        
        try:
            processor = self.supported_types[f'.{file_type}']
            return processor(file_path, filename)
        except Exception as e:
            return {
                'type': file_type,
                'content': '',
                'metadata': {'error': str(e)}
            }
    
    def _process_pdf(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Extract text and metadata from PDF"""
        content = []
        metadata = {'pages': 0, 'sections': []}
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata['pages'] = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        content.append({
                            'page': page_num + 1,
                            'text': page_text,
                            'type': 'text'
                        })
                        
                        # Extract potential headings
                        lines = page_text.split('\n')
                        for line in lines:
                            if self._is_heading(line):
                                metadata['sections'].append({
                                    'page': page_num + 1,
                                    'text': line.strip(),
                                    'type': 'heading'
                                })
        
        except Exception as e:
            metadata['error'] = str(e)
        
        return {
            'type': 'pdf',
            'content': content,
            'metadata': metadata,
            'filename': filename
        }
    
    def _process_pptx(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Extract content from PowerPoint presentation"""
        content = []
        metadata = {'slides': 0, 'sections': []}
        
        try:
            prs = Presentation(file_path)
            metadata['slides'] = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_content = {
                    'slide': slide_num + 1,
                    'title': '',
                    'text': [],
                    'shapes': [],
                    'type': 'slide'
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if shape == slide.shapes[0] and len(text) < 100:  # Likely a title
                            slide_content['title'] = text
                            metadata['sections'].append({
                                'slide': slide_num + 1,
                                'text': text,
                                'type': 'title'
                            })
                        else:
                            slide_content['text'].append(text)
                    
                    if hasattr(shape, "shape_type"):
                        slide_content['shapes'].append({
                            'type': str(shape.shape_type),
                            'text': getattr(shape, 'text', '')
                        })
                
                content.append(slide_content)
        
        except Exception as e:
            metadata['error'] = str(e)
        
        return {
            'type': 'pptx',
            'content': content,
            'metadata': metadata,
            'filename': filename
        }
    
    def _process_docx(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Extract content from Word document"""
        content = []
        metadata = {'paragraphs': 0, 'sections': []}
        
        try:
            doc = Document(file_path)
            metadata['paragraphs'] = len(doc.paragraphs)
            
            for para_num, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    para_content = {
                        'paragraph': para_num + 1,
                        'text': paragraph.text,
                        'style': paragraph.style.name if paragraph.style else 'Normal',
                        'type': 'paragraph'
                    }
                    content.append(para_content)
                    
                    # Check if it's a heading
                    if paragraph.style.name.startswith('Heading'):
                        metadata['sections'].append({
                            'paragraph': para_num + 1,
                            'text': paragraph.text,
                            'level': paragraph.style.name,
                            'type': 'heading'
                        })
        
        except Exception as e:
            metadata['error'] = str(e)
        
        return {
            'type': 'docx',
            'content': content,
            'metadata': metadata,
            'filename': filename
        }
    
    def _process_xlsx(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Extract content from Excel spreadsheet"""
        content = []
        metadata = {'sheets': 0, 'cells': 0}
        
        try:
            workbook = openpyxl.load_workbook(file_path)
            metadata['sheets'] = len(workbook.sheetnames)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = {
                    'sheet': sheet_name,
                    'data': [],
                    'type': 'sheet'
                }
                
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        sheet_content['data'].append(list(row))
                        metadata['cells'] += len([cell for cell in row if cell is not None])
                
                content.append(sheet_content)
        
        except Exception as e:
            metadata['error'] = str(e)
        
        return {
            'type': 'xlsx',
            'content': content,
            'metadata': metadata,
            'filename': filename
        }
    
    def _process_image(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Extract metadata from image"""
        try:
            with Image.open(file_path) as img:
                metadata = {
                    'width': img.width,
                    'height': img.height,
                    'format': img.format,
                    'mode': img.mode
                }
            
            return {
                'type': 'image',
                'content': f"Image file: {filename}",
                'metadata': metadata,
                'filename': filename
            }
        
        except Exception as e:
            return {
                'type': 'image',
                'content': '',
                'metadata': {'error': str(e)},
                'filename': filename
            }
    
    def _process_audio(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Extract text from audio using speech recognition"""
        content = []
        metadata = {'duration': 0, 'transcript': ''}
        
        try:
            # Convert audio to wav if needed
            audio = AudioSegment.from_file(file_path)
            metadata['duration'] = len(audio) / 1000.0  # Duration in seconds
            
            # Convert to wav for speech recognition
            wav_path = file_path.replace(os.path.splitext(file_path)[1], '.wav')
            audio.export(wav_path, format="wav")
            
            # Perform speech recognition
            r = sr.Recognizer()
            with sr.AudioFile(wav_path) as source:
                audio_data = r.record(source)
                try:
                    # Try Google Speech Recognition first
                    transcript = r.recognize_google(audio_data)
                except sr.UnknownValueError:
                    transcript = "Could not understand audio"
                except sr.RequestError as e:
                    transcript = f"Speech recognition error: {str(e)}"
                except Exception as e:
                    transcript = f"Audio processing error: {str(e)}"
                
                metadata['transcript'] = transcript
                content.append({
                    'text': transcript,
                    'type': 'transcript'
                })
            
            # Clean up temporary file
            os.remove(wav_path)
        
        except Exception as e:
            metadata['error'] = str(e)
        
        return {
            'type': 'audio',
            'content': content,
            'metadata': metadata,
            'filename': filename
        }
    
    def _process_video(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Extract frames and basic info from video"""
        content = []
        metadata = {'duration': 0, 'fps': 0, 'frames': 0}
        
        try:
            cap = cv2.VideoCapture(file_path)
            metadata['fps'] = cap.get(cv2.CAP_PROP_FPS)
            metadata['frames'] = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            metadata['duration'] = metadata['frames'] / metadata['fps'] if metadata['fps'] > 0 else 0
            
            # Extract key frames (every 30th frame)
            frame_count = 0
            while True:
                ret, frame = cap.read()
                if not ret:
                    break
                
                if frame_count % 30 == 0:  # Every 30th frame
                    content.append({
                        'frame': frame_count,
                        'timestamp': frame_count / metadata['fps'],
                        'type': 'frame'
                    })
                
                frame_count += 1
            
            cap.release()
        
        except Exception as e:
            metadata['error'] = str(e)
        
        return {
            'type': 'video',
            'content': content,
            'metadata': metadata,
            'filename': filename
        }
    
    def _is_heading(self, text: str) -> bool:
        """Check if text is likely a heading"""
        text = text.strip()
        if len(text) < 3 or len(text) > 100:
            return False
        
        # Check for common heading patterns
        heading_patterns = [
            r'^[A-Z][A-Z\s]+$',  # All caps
            r'^\d+\.?\s+[A-Z]',  # Numbered
            r'^[A-Z][a-z]+(\s+[A-Z][a-z]+)*$',  # Title case
        ]
        
        return any(re.match(pattern, text) for pattern in heading_patterns)

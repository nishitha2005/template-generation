import os
from datetime import datetime
from typing import Dict, Any, List
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json

class OutputGenerator:
    """Generate outputs in various formats (DOCX, PDF, PPTX)"""
    
    def __init__(self):
        self.supported_formats = ['docx', 'pdf', 'pptx']
    
    def generate_output(self, content: Dict[str, Any], template: Dict[str, Any], 
                      format: str, output_dir: str) -> str:
        """Generate output file in specified format"""
        
        if format not in self.supported_formats:
            raise ValueError(f"Unsupported format: {format}")
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"generated_content_{timestamp}.{format}"
        output_path = os.path.join(output_dir, filename)
        
        if format == 'docx':
            return self._generate_docx(content, template, output_path)
        elif format == 'pdf':
            return self._generate_pdf(content, template, output_path)
        elif format == 'pptx':
            return self._generate_pptx(content, template, output_path)
    
    def _generate_docx(self, content: Dict[str, Any], template: Dict[str, Any], output_path: str) -> str:
        """Generate DOCX file"""
        doc = Document()
        
        # Set up document formatting
        formatting = template.get('formatting', {})
        style = template.get('style', {})
        
        # Add title
        title = doc.add_heading(template.get('metadata', {}).get('name', 'Generated Document'), 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add metadata
        metadata = content.get('metadata', {})
        if metadata:
            doc.add_paragraph(f"Generated: {metadata.get('generated_at', 'Unknown')}")
            doc.add_paragraph(f"Sources: {', '.join(metadata.get('sources_used', []))}")
            doc.add_paragraph("")  # Empty line
        
        # Add sections in order
        sections = template.get('structure', {}).get('sections', [])
        for section_template in sections:
            section_id = section_template['id']
            section_title = section_template['title']
            
            if section_id in content.get('sections', {}):
                section_content = content['sections'][section_id]
                
                # Add section heading
                heading = doc.add_heading(section_title, level=1)
                
                # Add content based on type
                content_text = section_content.get('content', '')
                content_type = section_template.get('content_type', 'text')
                
                if content_type == 'list' and isinstance(content_text, list):
                    for item in content_text:
                        doc.add_paragraph(f"• {item}", style='List Bullet')
                else:
                    # Add as paragraph
                    para = doc.add_paragraph(str(content_text))
                
                # Add citations if any
                citations = section_content.get('citations', [])
                if citations:
                    doc.add_paragraph("Sources:", style='Heading 3')
                    for citation in citations:
                        doc.add_paragraph(f"• {citation.get('full_citation', '')}", style='List Bullet')
                
                doc.add_paragraph("")  # Empty line between sections
        
        # Save document
        doc.save(output_path)
        return output_path
    
    def _generate_pdf(self, content: Dict[str, Any], template: Dict[str, Any], output_path: str) -> str:
        """Generate PDF file"""
        doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Create custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=1  # Center alignment
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=14,
            spaceAfter=12,
            spaceBefore=20
        )
        
        # Add title
        title = template.get('metadata', {}).get('name', 'Generated Document')
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 20))
        
        # Add metadata
        metadata = content.get('metadata', {})
        if metadata:
            story.append(Paragraph(f"Generated: {metadata.get('generated_at', 'Unknown')}", styles['Normal']))
            story.append(Paragraph(f"Sources: {', '.join(metadata.get('sources_used', []))}", styles['Normal']))
            story.append(Spacer(1, 20))
        
        # Add sections
        sections = template.get('structure', {}).get('sections', [])
        for section_template in sections:
            section_id = section_template['id']
            section_title = section_template['title']
            
            if section_id in content.get('sections', {}):
                section_content = content['sections'][section_id]
                
                # Add section heading
                story.append(Paragraph(section_title, heading_style))
                
                # Add content
                content_text = section_content.get('content', '')
                content_type = section_template.get('content_type', 'text')
                
                if content_type == 'list' and isinstance(content_text, list):
                    for item in content_text:
                        story.append(Paragraph(f"• {item}", styles['Normal']))
                else:
                    story.append(Paragraph(str(content_text), styles['Normal']))
                
                # Add citations
                citations = section_content.get('citations', [])
                if citations:
                    story.append(Paragraph("Sources:", styles['Heading3']))
                    for citation in citations:
                        story.append(Paragraph(f"• {citation.get('full_citation', '')}", styles['Normal']))
                
                story.append(Spacer(1, 12))
        
        # Build PDF
        doc.build(story)
        return output_path
    
    def _generate_pptx(self, content: Dict[str, Any], template: Dict[str, Any], output_path: str) -> str:
        """Generate PPTX file"""
        prs = Presentation()
        
        # Set slide size to widescreen
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = template.get('metadata', {}).get('name', 'Generated Presentation')
        subtitle.text = f"Generated: {content.get('metadata', {}).get('generated_at', 'Unknown')}"
        
        # Add content slides
        sections = template.get('structure', {}).get('sections', [])
        for section_template in sections:
            section_id = section_template['id']
            section_title = section_template['title']
            
            if section_id in content.get('sections', {}):
                section_content = content['sections'][section_id]
                
                # Create new slide
                slide_layout = prs.slide_layouts[1]  # Title and content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title = slide.shapes.title
                title.text = section_title
                
                # Add content
                content_text = section_content.get('content', '')
                content_type = section_template.get('content_type', 'text')
                
                # Get content placeholder
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                if content_type == 'list' and isinstance(content_text, list):
                    for i, item in enumerate(content_text):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = f"• {item}"
                        p.font.size = Pt(18)
                else:
                    p = text_frame.paragraphs[0]
                    p.text = str(content_text)
                    p.font.size = Pt(18)
                
                # Add citations if any
                citations = section_content.get('citations', [])
                if citations:
                    # Add another slide for citations if there are many
                    if len(citations) > 3:
                        citation_slide = prs.slides.add_slide(prs.slide_layouts[1])
                        citation_slide.shapes.title.text = f"{section_title} - Sources"
                        
                        citation_text_frame = citation_slide.placeholders[1].text_frame
                        citation_text_frame.clear()
                        
                        for citation in citations:
                            p = citation_text_frame.add_paragraph()
                            p.text = f"• {citation.get('full_citation', '')}"
                            p.font.size = Pt(14)
        
        # Save presentation
        prs.save(output_path)
        return output_path
    
    def generate_summary_report(self, content: Dict[str, Any], template: Dict[str, Any]) -> Dict[str, Any]:
        """Generate a summary report of the generated content"""
        
        summary = {
            "document_info": {
                "title": template.get('metadata', {}).get('name', 'Unknown'),
                "generated_at": content.get('metadata', {}).get('generated_at', 'Unknown'),
                "template_version": template.get('metadata', {}).get('version', '1.0'),
                "sources_used": content.get('metadata', {}).get('sources_used', [])
            },
            "content_stats": {
                "total_sections": len(content.get('sections', {})),
                "total_word_count": sum(
                    section.get('word_count', 0) 
                    for section in content.get('sections', {}).values()
                ),
                "total_citations": sum(
                    len(section.get('citations', [])) 
                    for section in content.get('sections', {}).values()
                )
            },
            "sections": []
        }
        
        # Add section summaries
        for section_id, section in content.get('sections', {}).items():
            section_summary = {
                "id": section_id,
                "title": section.get('title', ''),
                "word_count": section.get('word_count', 0),
                "citation_count": len(section.get('citations', [])),
                "content_preview": str(section.get('content', ''))[:200] + "..." if len(str(section.get('content', ''))) > 200 else str(section.get('content', ''))
            }
            summary["sections"].append(section_summary)
        
        return summary
